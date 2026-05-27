"""오프라인 모드: python-pptx로 .pptx 파일을 직접 생성/수정.

PowerPoint가 켜져있지 않아도 동작 — 배치 생성, CI, 헤드리스 환경용.
라이브 모드와 같은 API 인터페이스를 가능한 한 맞춘다.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from pptmaker import themes
from pptmaker.design_tokens import TOKENS

LAYOUTS_DIR = themes.LAYOUTS_DIR
DEFAULT_TEMPLATE = themes.template_path(themes.DEFAULT_THEME)


class OfflineBuilder:
    """python-pptx 기반 오프라인 PPT 빌더.

    Args:
        theme: 'theme1' | 'theme2'. None이면 DEFAULT_THEME 사용.
        template_path: theme보다 우선. 임의 경로로 템플릿 강제 지정 (고급).
    """

    def __init__(
        self,
        theme: str | None = None,
        template_path: Path | str | None = None,
    ):
        if template_path is not None:
            self._path = Path(template_path)
        else:
            self._path = themes.template_path(theme)
        if not self._path.exists():
            raise FileNotFoundError(f"Template not found: {self._path}")
        self._theme = theme or themes.DEFAULT_THEME
        self._prs = Presentation(str(self._path))

    @property
    def theme(self) -> str:
        return self._theme

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides)

    def _layout_by_name(self, *needles: str):
        for layout in self._prs.slide_layouts:
            if any(n in layout.name for n in needles):
                return layout
        return self._prs.slide_layouts[1]  # fallback: 제목 및 내용

    def add_body_slide(self, title: str, body_lines: list[str]) -> int:
        layout = self._layout_by_name("제목 및 내용", "Title and Content")
        slide = self._prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        for ph in slide.placeholders:
            if ph == slide.shapes.title:
                continue
            if ph.has_text_frame:
                tf = ph.text_frame
                tf.text = body_lines[0] if body_lines else ""
                for line in body_lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.name = TOKENS.fonts.minor
                    p.font.size = Pt(18)
                break
        return len(self._prs.slides)

    def add_title_only_slide(self, title: str) -> int:
        layout = self._layout_by_name("제목만", "Title Only")
        slide = self._prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
        return len(self._prs.slides)

    def save(self, path: Path | str) -> Path:
        out = Path(path).resolve()
        out.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(out))
        return out
