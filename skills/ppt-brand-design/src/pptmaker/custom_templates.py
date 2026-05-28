"""사용자 커스텀 레이아웃 — Phase 1 인프라.

사용자가 user_layouts/<theme>_user.pptx 에 슬라이드를 디자인해 두고
슬라이드 노트에 `pptmaker:custom name=<key> placeholders=<csv>` 메타데이터를
넣으면, Claude/코드가 그 슬라이드를 복제 + {{placeholder}} 치환해서 활용할 수 있다.

manifest.json (자동/수동 갱신)은 디렉토리·메타 캐시 역할로 Claude가 빠르게
사용 가능 목록을 파악할 수 있게 한다 (없어도 .pptx 스캔으로 대체 가능).

핵심 원칙:
- layouts/ (사내 배포본) 은 절대 안 건드림 — user_layouts/ 만 활용
- git pull 시 사용자 추가분 보존 (user_layouts/ 는 .gitignore)
- 사용자 입장에선 "theme1에 추가했다" 처럼 transparent
"""
from __future__ import annotations

import json
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterator

from pptmaker import themes


META_LINE_RE = re.compile(
    r"pptmaker:custom\s+name=(\S+)\s+placeholders=([^\s]*)"
)
PLACEHOLDER_RE = re.compile(r"\{\{\s*([^}]+?)\s*\}\}")


@dataclass
class CustomTemplate:
    """사용자 커스텀 레이아웃 한 항목."""
    name: str
    theme: str
    source_file: str      # user_layouts/ 안의 .pptx 파일명 (예: "theme1_user.pptx")
    source_slide_index: int  # 1-based, 그 .pptx 안의 슬라이드 위치
    placeholders: list[str] = field(default_factory=list)
    description: str = ""

    def to_dict(self) -> dict:
        return {
            "name": self.name,
            "theme": self.theme,
            "source_file": self.source_file,
            "source_slide_index": self.source_slide_index,
            "placeholders": self.placeholders,
            "description": self.description,
        }

    @classmethod
    def from_dict(cls, d: dict) -> "CustomTemplate":
        return cls(
            name=d["name"],
            theme=d["theme"],
            source_file=d["source_file"],
            source_slide_index=int(d["source_slide_index"]),
            placeholders=list(d.get("placeholders") or []),
            description=str(d.get("description") or ""),
        )


# ── manifest.json 로드/저장 ───────────────────────────────────────────
def load_manifest() -> list[CustomTemplate]:
    """user_layouts/manifest.json 읽어 CustomTemplate 리스트 반환. 없으면 빈 리스트."""
    p = themes.user_manifest_path()
    if not p.exists():
        return []
    try:
        data = json.loads(p.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return []
    items = data.get("templates", []) if isinstance(data, dict) else []
    return [CustomTemplate.from_dict(d) for d in items if isinstance(d, dict)]


def save_manifest(templates: list[CustomTemplate]) -> None:
    """manifest.json 저장."""
    themes.USER_LAYOUTS_DIR.mkdir(parents=True, exist_ok=True)
    data = {"templates": [t.to_dict() for t in templates]}
    themes.user_manifest_path().write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8"
    )


# ── .pptx 스캔으로 메타데이터 추출 (manifest 없거나 sync용) ────────────
def _read_notes_text(slide) -> str:
    """python-pptx Slide 의 notes 텍스트 반환 (없으면 빈 문자열)."""
    try:
        if not slide.has_notes_slide:
            return ""
        return slide.notes_slide.notes_text_frame.text or ""
    except Exception:
        return ""


def parse_meta_from_notes(notes_text: str) -> tuple[str, list[str]] | None:
    """슬라이드 노트에서 pptmaker:custom 메타 추출.

    형식: "pptmaker:custom name=<id> placeholders=key1,key2,key3"
    반환: (name, [placeholder, ...]) 또는 None.
    """
    m = META_LINE_RE.search(notes_text or "")
    if not m:
        return None
    name = m.group(1)
    csv = m.group(2).strip()
    placeholders = [p.strip() for p in csv.split(",") if p.strip()] if csv else []
    return name, placeholders


def scan_user_pptx(theme: str) -> list[CustomTemplate]:
    """user_layouts/<theme>_user.pptx 를 python-pptx로 스캔.

    각 슬라이드의 노트를 검사하고 메타데이터가 있으면 CustomTemplate 으로 등록.
    """
    from pptx import Presentation  # 지연 import (pywin32 없는 환경에서도 import OK)

    src = themes.user_template_path(theme)
    if not src.exists():
        return []
    prs = Presentation(str(src))
    out: list[CustomTemplate] = []
    for i, slide in enumerate(prs.slides, start=1):
        notes = _read_notes_text(slide)
        meta = parse_meta_from_notes(notes)
        if meta is None:
            continue
        name, placeholders = meta
        out.append(CustomTemplate(
            name=name,
            theme=theme,
            source_file=src.name,
            source_slide_index=i,
            placeholders=placeholders,
        ))
    return out


def list_templates(theme: str | None = None) -> list[CustomTemplate]:
    """사용 가능한 사용자 커스텀 레이아웃 목록.

    우선순위:
      1) manifest.json 이 있으면 그걸 사용 (빠름)
      2) 없으면 user_layouts/<theme>_user.pptx 를 스캔

    theme 지정 시 그 테마만 필터.
    """
    templates = load_manifest()
    if not templates:
        # manifest 없음 — 두 테마 모두 스캔
        for t in themes.VALID_THEMES:
            templates.extend(scan_user_pptx(t))
    if theme is None:
        return templates
    return [t for t in templates if t.theme == theme]


def find_template(name: str, theme: str | None = None) -> CustomTemplate | None:
    """이름으로 템플릿 검색."""
    for t in list_templates(theme):
        if t.name == name:
            return t
    return None


# ── {{key}} placeholder 치환 (COM 슬라이드) ───────────────────────────
def replace_placeholders_com(slide, mapping: dict[str, str]) -> int:
    """COM 슬라이드 도형들을 순회하며 {{key}} 패턴을 mapping[key] 로 치환.

    Returns: 치환된 항목 개수.
    """
    count = 0
    for shape in slide.Shapes:
        try:
            if shape.HasTextFrame == -1:  # MSO_TRUE
                tr = shape.TextFrame.TextRange
                original = tr.Text
                if not original:
                    continue
                new_text = PLACEHOLDER_RE.sub(
                    lambda m: str(mapping.get(m.group(1).strip(), m.group(0))),
                    original,
                )
                if new_text != original:
                    tr.Text = new_text
                    count += 1
        except Exception:
            continue
    return count


def detect_placeholders_in_com_slide(slide) -> list[str]:
    """COM 슬라이드의 모든 텍스트박스를 스캔해서 {{key}} 키들을 수집.

    중복 제거 + 등장 순서 보존.
    """
    seen: dict[str, None] = {}  # ordered set
    for shape in slide.Shapes:
        try:
            if shape.HasTextFrame == -1:
                text = shape.TextFrame.TextRange.Text or ""
                for m in PLACEHOLDER_RE.finditer(text):
                    key = m.group(1).strip()
                    if key not in seen:
                        seen[key] = None
        except Exception:
            continue
    return list(seen.keys())


def add_to_manifest(template: CustomTemplate, *, overwrite: bool = False) -> None:
    """manifest.json 에 항목 추가.

    중복 name 발견 시:
      - overwrite=True → 기존을 새 것으로 교체
      - overwrite=False → ValueError
    """
    items = load_manifest()
    existing_idx = next((i for i, t in enumerate(items) if t.name == template.name), None)
    if existing_idx is not None:
        if not overwrite:
            raise ValueError(
                f"이미 등록된 이름: '{template.name}'. overwrite=True 로 덮어쓰거나 다른 이름을 사용하세요."
            )
        items[existing_idx] = template
    else:
        items.append(template)
    save_manifest(items)


def create_user_pptx_if_missing(theme: str) -> Path:
    """user_layouts/<theme>_user.pptx 가 없으면 layouts/<theme>.pptx 의 chrome만
    복사해서 빈 상태로 생성. 있으면 그대로 반환.

    빈 상태 = 데모 본문 슬라이드는 모두 제거하고 cover + closing 만 남김.
    이후 register_layout() 호출이 슬라이드를 append 한다.
    """
    target = themes.user_template_path(theme)
    if target.exists():
        return target

    # layouts/<theme>.pptx 복사
    from pptx import Presentation
    import shutil

    src = themes.template_path(theme)
    themes.USER_LAYOUTS_DIR.mkdir(parents=True, exist_ok=True)
    shutil.copy2(str(src), str(target))

    # 데모 본문 슬라이드 (2..N-1) 제거 — cover(1) + closing(last) 만 남김
    prs = Presentation(str(target))
    sldIdLst = prs.slides._sldIdLst
    slide_ids = list(sldIdLst)
    if len(slide_ids) >= 3:
        # 첫(cover)과 마지막(closing) 외 다 제거
        for sid in slide_ids[1:-1]:
            sldIdLst.remove(sid)
            try:
                prs.part.drop_rel(sid.rId)
            except Exception:
                pass
        prs.save(str(target))

    return target


def build_meta_line(name: str, placeholders: list[str], description: str = "") -> str:
    """슬라이드 노트에 들어갈 메타 한 줄 생성."""
    csv = ",".join(placeholders)
    base = f"pptmaker:custom name={name} placeholders={csv}"
    if description:
        return f"{base}\n{description}"
    return base
