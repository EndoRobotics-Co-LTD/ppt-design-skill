"""Reference .pptx 두 개를 분석해 design_tokens.json을 생성한다.

추출 대상:
- 슬라이드 크기/비율
- 슬라이드 마스터 + 레이아웃 목록
- 테마 컬러 (XML 직접 파싱)
- 사용된 폰트 패밀리/사이즈
- 각 슬라이드의 placeholder 구조
- 이미지/도형 개수 요약
"""
from __future__ import annotations

import json
import re
from collections import Counter
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
REFERENCE_DIR = ROOT / "Reference"
OUTPUT = ROOT / "design_tokens.json"

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def emu_to_cm(emu: int) -> float:
    return round(Emu(emu).cm, 2)


def extract_theme_colors(pptx_path: Path) -> dict:
    """theme1.xml의 색상 스킴을 그대로 추출."""
    import zipfile

    colors: dict[str, str] = {}
    with zipfile.ZipFile(pptx_path) as zf:
        theme_files = [n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
        if not theme_files:
            return colors
        with zf.open(theme_files[0]) as f:
            tree = ET.parse(f)
        root = tree.getroot()
        scheme = root.find(".//a:clrScheme", NS)
        if scheme is None:
            return colors
        for elem in scheme:
            tag = elem.tag.split("}")[-1]
            srgb = elem.find("a:srgbClr", NS)
            sys = elem.find("a:sysClr", NS)
            if srgb is not None:
                colors[tag] = "#" + srgb.get("val", "").upper()
            elif sys is not None:
                colors[tag] = "sys:" + sys.get("lastClr", sys.get("val", ""))
        return colors


def extract_theme_fonts(pptx_path: Path) -> dict:
    """theme1.xml의 majorFont/minorFont (제목/본문 폰트) 추출."""
    import zipfile

    fonts: dict[str, str] = {}
    with zipfile.ZipFile(pptx_path) as zf:
        theme_files = [n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
        if not theme_files:
            return fonts
        with zf.open(theme_files[0]) as f:
            tree = ET.parse(f)
        root = tree.getroot()
        font_scheme = root.find(".//a:fontScheme", NS)
        if font_scheme is None:
            return fonts
        for role in ("majorFont", "minorFont"):
            node = font_scheme.find(f"a:{role}", NS)
            if node is None:
                continue
            latin = node.find("a:latin", NS)
            ea = node.find("a:ea", NS)
            fonts[role] = {
                "latin": latin.get("typeface") if latin is not None else None,
                "ea": ea.get("typeface") if ea is not None else None,
            }
    return fonts


def collect_runtime_fonts(prs: Presentation) -> dict:
    """실제 슬라이드에서 사용된 폰트와 크기 분포."""
    families: Counter[str] = Counter()
    sizes: Counter[float] = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        families[run.font.name] += 1
                    if run.font.size:
                        sizes[run.font.size.pt] += 1
    return {
        "families": families.most_common(),
        "sizes_pt": sizes.most_common(),
    }


def collect_layouts(prs: Presentation) -> list[dict]:
    layouts = []
    for layout in prs.slide_layouts:
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type).split(".")[-1] if ph.placeholder_format.type else None,
                "name": ph.name,
            })
        layouts.append({
            "name": layout.name,
            "placeholders": placeholders,
        })
    return layouts


def summarize_slides(prs: Presentation) -> list[dict]:
    summary = []
    for i, slide in enumerate(prs.slides, start=1):
        shape_counter: Counter[str] = Counter()
        text_snippets: list[str] = []
        for shape in slide.shapes:
            shape_counter[shape.shape_type.__class__.__name__ if shape.shape_type is None else str(shape.shape_type).split(".")[-1]] += 1
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    text_snippets.append(re.sub(r"\s+", " ", txt)[:120])
        summary.append({
            "index": i,
            "layout": slide.slide_layout.name,
            "shape_counts": dict(shape_counter),
            "text_preview": text_snippets[:3],
        })
    return summary


def analyze(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))
    return {
        "file": pptx_path.name,
        "slide_size": {
            "width_cm": emu_to_cm(prs.slide_width),
            "height_cm": emu_to_cm(prs.slide_height),
            "aspect": round(prs.slide_width / prs.slide_height, 4),
        },
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "theme_colors": extract_theme_colors(pptx_path),
        "theme_fonts": extract_theme_fonts(pptx_path),
        "runtime_fonts": collect_runtime_fonts(prs),
        "layouts": collect_layouts(prs),
        "slides": summarize_slides(prs),
    }


def main() -> None:
    files = sorted(REFERENCE_DIR.glob("*.pptx"))
    if not files:
        raise SystemExit(f"No .pptx files found in {REFERENCE_DIR}")
    report = {"templates": [analyze(p) for p in files]}
    OUTPUT.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"OK -> {OUTPUT} ({OUTPUT.stat().st_size} bytes)")
    for tpl in report["templates"]:
        print(f"\n=== {tpl['file']} ===")
        print(f"  slides: {tpl['slide_count']}, layouts: {tpl['layout_count']}")
        print(f"  size: {tpl['slide_size']}")
        print(f"  theme_colors: {tpl['theme_colors']}")
        print(f"  theme_fonts: {tpl['theme_fonts']}")
        print(f"  top runtime fonts: {tpl['runtime_fonts']['families'][:5]}")
        print(f"  top runtime sizes: {tpl['runtime_fonts']['sizes_pt'][:5]}")


if __name__ == "__main__":
    main()
