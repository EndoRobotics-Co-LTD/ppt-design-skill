"""Reference 템플릿의 슬라이드 마스터 + 각 슬라이드의 도형을 상세 조사.

목적:
- 마스터에 어떤 도형이 있는지 (로고, 장식, 페이지 번호 등)
- 각 슬라이드의 실제 도형 위치/크기/텍스트
- 사용된 그림 파일 (로고 추출 대상)
"""
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "src"))

from pptx import Presentation
from pptx.util import Emu

REFERENCE_DIR = ROOT / "Reference"


def cm(emu):
    return round(Emu(emu).cm, 2) if emu else 0.0


def describe_shape(shape, indent=2):
    pre = " " * indent
    typ = str(shape.shape_type).split(".")[-1] if shape.shape_type else "?"
    info = f"{pre}- {typ} | name='{shape.name}' | pos=({cm(shape.left)},{cm(shape.top)}) size=({cm(shape.width)},{cm(shape.height)})"
    extras = []
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.strip().replace("\n", " | ")[:80]
        extras.append(f"text='{text}'")
    if shape.shape_type == 13:  # PICTURE
        try:
            img = shape.image
            extras.append(f"image_ext={img.ext} size={len(img.blob)//1024}KB")
        except Exception:
            pass
    if extras:
        info += " | " + " ".join(extras)
    print(info)


def inspect_pptx(path: Path):
    print(f"\n{'=' * 70}")
    print(f"FILE: {path.name}")
    print(f"{'=' * 70}")
    prs = Presentation(str(path))

    # 슬라이드 마스터 (보통 1개)
    print(f"\n[Slide Masters: {len(prs.slide_masters)}]")
    for mi, master in enumerate(prs.slide_masters, 1):
        print(f"\n  Master {mi}: shapes on master itself")
        for shape in master.shapes:
            describe_shape(shape, indent=4)
        print(f"\n  Master {mi}: {len(master.slide_layouts)} layouts")
        for li, layout in enumerate(master.slide_layouts, 1):
            print(f"    Layout {li}: '{layout.name}' ({len(layout.shapes)} shapes)")
            for shape in layout.shapes:
                describe_shape(shape, indent=6)

    # 실제 슬라이드들
    print(f"\n[Slides: {len(prs.slides)}]")
    for si, slide in enumerate(prs.slides, 1):
        print(f"\n  Slide {si} (layout: '{slide.slide_layout.name}', shapes: {len(slide.shapes)})")
        for shape in slide.shapes:
            describe_shape(shape, indent=4)


def extract_pictures(path: Path, out_dir: Path):
    """zip 내부의 media/* 파일들을 그대로 저장 — 로고 같은 그림 추출."""
    out_dir.mkdir(exist_ok=True, parents=True)
    with zipfile.ZipFile(path) as zf:
        media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
        for n in media:
            data = zf.read(n)
            target = out_dir / (path.stem.replace(" ", "_") + "__" + Path(n).name)
            target.write_bytes(data)
            print(f"  saved {target.name} ({len(data)//1024}KB)")
        return len(media)


def main():
    pptx_files = [p for p in sorted(REFERENCE_DIR.glob("*.pptx")) if not p.name.startswith("~$")]
    for pptx in pptx_files:
        inspect_pptx(pptx)

    print("\n\n=== 그림 파일 추출 ===")
    out_dir = ROOT / "Reference" / "_extracted_media"
    total = 0
    for pptx in pptx_files:
        print(f"\n{pptx.name}:")
        total += extract_pictures(pptx, out_dir)
    print(f"\n총 {total}개 그림 추출 → {out_dir}")


if __name__ == "__main__":
    main()
